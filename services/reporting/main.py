"""
CLISONIX REPORTING SERVICE - REAL DATA EXCEL
Excel i vërtetë me tabela të plota, jo fake, jo mock
Port: 8001
"""

import asyncio
import json
import logging
import os
import subprocess
import time
from datetime import datetime
from pathlib import Path
from typing import Any, Callable, Dict, List, Optional, Sequence, cast

import httpx
from fastapi import FastAPI, HTTPException, Query, Response
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse

# Setup logging
logging.basicConfig(level=logging.INFO, format='%(asctime)s | %(levelname)s | %(message)s')
logger = logging.getLogger("reporting-real")

app = FastAPI(
    title="Clisonix Reporting - Real Client Intelligence",
    description="Raporte, analiza dhe statistika reale për klientë dhe operacione production-grade.",
    version="5.0.0",
    docs_url="/docs"
)

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

BASE_DIR = Path(__file__).resolve().parent
DEFAULT_PROJECT_ROOT = BASE_DIR.parents[1] if len(BASE_DIR.parents) > 1 else BASE_DIR.parent
PROJECT_ROOT = Path(os.getenv("CLISONIX_PROJECT_ROOT", str(DEFAULT_PROJECT_ROOT))).resolve()
REPORTS_DIR = BASE_DIR / "reports"
REPORTS_DIR.mkdir(parents=True, exist_ok=True)

IGNORED_SCAN_DIRS = {
    ".git", ".hg", ".svn", ".next", "node_modules", "dist", "build",
    "__pycache__", ".venv", "venv", ".mypy_cache", ".pytest_cache"
}
MATERIALS_CACHE_TTL_SECONDS = int(os.getenv("REPORTING_MATERIALS_CACHE_TTL", "120"))
MATERIALS_SCAN_MAX_SECONDS = float(os.getenv("REPORTING_MATERIALS_SCAN_MAX_SECONDS", "2.5"))
MATERIALS_SCAN_MAX_FILES = int(os.getenv("REPORTING_MATERIALS_SCAN_MAX_FILES", "6000"))
UPSTREAM_CACHE_TTL_SECONDS = int(os.getenv("REPORTING_UPSTREAM_CACHE_TTL", "30"))
UPSTREAM_REQUEST_TIMEOUT_SECONDS = float(os.getenv("REPORTING_UPSTREAM_REQUEST_TIMEOUT", "1.75"))
BLOCKING_CALL_TIMEOUT_SECONDS = float(os.getenv("REPORTING_BLOCKING_CALL_TIMEOUT", "1.5"))
MAX_HISTORY_POINTS = int(os.getenv("REPORTING_MAX_HISTORY_POINTS", "720"))

MATERIALS_CACHE: Dict[str, Any] = {"fetched_at": 0.0, "data": None}
UPSTREAM_CACHE: Dict[str, Any] = {"fetched_at": 0.0, "data": None}
DOCKER_CONTAINERS_CACHE: Dict[str, Any] = {"fetched_at": 0.0, "data": None}
DOCKER_STATS_CACHE: Dict[str, Any] = {"fetched_at": 0.0, "data": None}
DOCKER_CACHE_TTL_SECONDS = int(os.getenv("REPORTING_DOCKER_CACHE_TTL", "20"))
METRIC_HISTORY: List[Dict[str, Any]] = []

# Import libraries
try:
    from openpyxl import Workbook  # type: ignore[import-not-found]
    from openpyxl.styles import (  # type: ignore[import-not-found]
        Alignment,
        Border,
        Font,
        PatternFill,
        Side,
    )
    from openpyxl.utils import get_column_letter  # type: ignore[import-not-found]
    from openpyxl.worksheet.table import (  # type: ignore[import-not-found]
        Table,
        TableStyleInfo,
    )
    EXCEL_AVAILABLE = True
except ImportError:
    EXCEL_AVAILABLE = False
    Workbook = None
    Font = PatternFill = Alignment = Border = Side = get_column_letter = None
    Table = None
    TableStyleInfo = None
    logger.error("openpyxl not installed")

try:
    import pptx  # type: ignore[import-not-found]
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    pptx = None  # type: ignore[assignment]

try:
    import psutil  # type: ignore[import-not-found]
    PSUTIL_AVAILABLE = True
except ImportError:
    PSUTIL_AVAILABLE = False
    psutil = None

# Docker SDK
try:
    import docker  # type: ignore[import-not-found]
except ImportError:
    docker = None

DOCKER_SDK_AVAILABLE = docker is not None
docker_client = None


def get_docker_client():
    """Create or reuse a Docker client lazily so startup races do not permanently disable container reporting."""
    global docker_client, DOCKER_SDK_AVAILABLE

    if docker is None:
        DOCKER_SDK_AVAILABLE = False
        return None

    if docker_client is not None:
        return docker_client

    try:
        docker_client = docker.from_env()
        DOCKER_SDK_AVAILABLE = True
        return docker_client
    except Exception as e:
        DOCKER_SDK_AVAILABLE = False
        docker_client = None
        logger.warning(f"Docker client unavailable: {e}")
        return None


def get_docker_containers_real(force_refresh: bool = False):
    """Merr container-ët real duke preferuar Docker CLI për përgjigje më të shpejta."""
    now = time.time()
    cached = DOCKER_CONTAINERS_CACHE.get("data")
    if not force_refresh and cached and (now - float(DOCKER_CONTAINERS_CACHE.get("fetched_at", 0.0))) < DOCKER_CACHE_TTL_SECONDS:
        return cached

    containers = []

    # Prefero Docker CLI sepse është dukshëm më i shpejtë se inspect-i për çdo container.
    try:
        result = subprocess.run(
            ["docker", "ps", "--format", "{{.Names}}|{{.Status}}|{{.Ports}}|{{.Image}}|{{.ID}}"],
            capture_output=True, text=True, timeout=8
        )
        if result.returncode == 0:
            for line in result.stdout.strip().splitlines():
                if not line:
                    continue
                parts = line.split('|')
                if len(parts) < 5:
                    continue
                status_text = parts[1]
                normalized_status = status_text.lower()
                containers.append({
                    "name": parts[0],
                    "status": status_text,
                    "ports": parts[2][:60] if parts[2] else "-",
                    "image": parts[3],
                    "container_id": parts[4][:12],
                    "healthy": ("healthy" in normalized_status) or normalized_status.startswith("up") or ("running" in normalized_status),
                    "uptime": status_text,
                })
            DOCKER_CONTAINERS_CACHE["fetched_at"] = now
            DOCKER_CONTAINERS_CACHE["data"] = containers
            return containers
    except Exception as e:
        logger.warning(f"Docker CLI container scan failed: {e}")

    client = get_docker_client()
    if client:
        try:
            for c in client.containers.list(sparse=True):
                status_text = str(getattr(c, "status", "unknown"))
                containers.append({
                    "name": getattr(c, "name", "unknown"),
                    "status": status_text,
                    "ports": "-",
                    "image": "docker-sdk",
                    "container_id": getattr(c, "short_id", "unknown"),
                    "healthy": status_text == "running",
                    "uptime": status_text,
                })
        except Exception as e:
            logger.error(f"Docker SDK container scan failed: {e}")

    DOCKER_CONTAINERS_CACHE["fetched_at"] = now
    DOCKER_CONTAINERS_CACHE["data"] = containers
    return containers


def get_docker_stats_real(force_refresh: bool = False):
    """Merr CPU/Memory stats real me cache të shkurtër për të shmangur timeout-et."""
    now = time.time()
    cached = DOCKER_STATS_CACHE.get("data")
    if not force_refresh and cached and (now - float(DOCKER_STATS_CACHE.get("fetched_at", 0.0))) < DOCKER_CACHE_TTL_SECONDS:
        return cached

    stats = []

    # Prefero një thirrje të vetme në Docker CLI.
    try:
        result = subprocess.run(
            ["docker", "stats", "--no-stream", "--format",
             "{{.Name}}|{{.CPUPerc}}|{{.MemUsage}}|{{.MemPerc}}|{{.NetIO}}|{{.BlockIO}}"],
            capture_output=True, text=True, timeout=10
        )
        if result.returncode == 0:
            for line in result.stdout.strip().splitlines():
                if not line:
                    continue
                parts = line.split('|')
                if len(parts) < 6:
                    continue
                stats.append({
                    "container": parts[0],
                    "cpu": parts[1],
                    "mem_usage": parts[2],
                    "mem_percent": parts[3],
                    "net_io": parts[4],
                    "block_io": parts[5],
                })
            DOCKER_STATS_CACHE["fetched_at"] = now
            DOCKER_STATS_CACHE["data"] = stats
            return stats
    except Exception as e:
        logger.warning(f"Docker CLI stats failed: {e}")

    client = get_docker_client()
    if client:
        try:
            for c in client.containers.list()[:12]:
                try:
                    stats_payload = cast(Dict[str, Any], c.stats(stream=False))
                    cpu_stats = cast(Dict[str, Any], stats_payload.get('cpu_stats', {}))
                    pre_cpu_stats = cast(Dict[str, Any], stats_payload.get('precpu_stats', {}))
                    cpu_usage = cast(Dict[str, Any], cpu_stats.get('cpu_usage', {}))
                    pre_cpu_usage = cast(Dict[str, Any], pre_cpu_stats.get('cpu_usage', {}))
                    memory_stats = cast(Dict[str, Any], stats_payload.get('memory_stats', {}))

                    cpu_delta = _safe_float(cpu_usage.get('total_usage')) - _safe_float(pre_cpu_usage.get('total_usage'))
                    system_delta = _safe_float(cpu_stats.get('system_cpu_usage')) - _safe_float(pre_cpu_stats.get('system_cpu_usage'))
                    cpu_percent = (cpu_delta / system_delta) * 100 if system_delta > 0 else 0

                    mem_usage = _safe_float(memory_stats.get('usage')) / (1024 * 1024)
                    mem_limit = max(_safe_float(memory_stats.get('limit'), 1.0) / (1024 * 1024), 1.0)
                    mem_percent = (mem_usage / mem_limit) * 100 if mem_limit > 0 else 0

                    stats.append({
                        "container": getattr(c, 'name', 'unknown'),
                        "cpu": f"{cpu_percent:.2f}%",
                        "mem_usage": f"{mem_usage:.1f}MiB / {mem_limit:.0f}MiB",
                        "mem_percent": f"{mem_percent:.2f}%",
                        "net_io": "N/A",
                        "block_io": "N/A",
                    })
                except Exception as e:
                    logger.error(f"Stats error for {getattr(c, 'name', 'unknown')}: {e}")
        except Exception as e:
            logger.error(f"Docker SDK stats error: {e}")

    DOCKER_STATS_CACHE["fetched_at"] = now
    DOCKER_STATS_CACHE["data"] = stats
    return stats


def _format_uptime(seconds: int) -> str:
    days, remainder = divmod(max(0, seconds), 86400)
    hours, remainder = divmod(remainder, 3600)
    minutes, _ = divmod(remainder, 60)
    return f"{days}d {hours}h {minutes}m"


def get_system_metrics_real() -> Dict[str, Any]:
    """Merr metrika REALE nga sistemi"""
    metrics: Dict[str, Any] = {
        "collector": "psutil" if PSUTIL_AVAILABLE and psutil is not None else "unavailable",
        "timestamp": datetime.now().isoformat(),
    }

    if PSUTIL_AVAILABLE and psutil is not None:
        metrics["cpu_percent"] = float(psutil.cpu_percent(interval=0.3))
        metrics["cpu_count"] = int(psutil.cpu_count() or 0)

        mem = psutil.virtual_memory()
        metrics["memory_total_gb"] = round(mem.total / (1024**3), 2)
        metrics["memory_used_gb"] = round(mem.used / (1024**3), 2)
        metrics["memory_percent"] = float(mem.percent)

        disk_target = str(PROJECT_ROOT) if PROJECT_ROOT.exists() else str(Path.cwd())
        disk = psutil.disk_usage(disk_target)
        metrics["disk_total_gb"] = round(disk.total / (1024**3), 2)
        metrics["disk_used_gb"] = round(disk.used / (1024**3), 2)
        metrics["disk_percent"] = round(float(disk.percent), 1)

        net = psutil.net_io_counters()
        metrics["net_sent_gb"] = round(net.bytes_sent / (1024**3), 2)
        metrics["net_recv_gb"] = round(net.bytes_recv / (1024**3), 2)

        boot_time = datetime.fromtimestamp(psutil.boot_time())
        uptime = datetime.now() - boot_time
        uptime_seconds = int(uptime.total_seconds())
        metrics["uptime_seconds"] = uptime_seconds
        metrics["uptime_hours"] = round(uptime_seconds / 3600, 1)
        metrics["uptime_formatted"] = _format_uptime(uptime_seconds)
    else:
        metrics.update({
            "cpu_percent": 0.0,
            "cpu_count": 0,
            "memory_total_gb": 0.0,
            "memory_used_gb": 0.0,
            "memory_percent": 0.0,
            "disk_total_gb": 0.0,
            "disk_used_gb": 0.0,
            "disk_percent": 0.0,
            "net_sent_gb": 0.0,
            "net_recv_gb": 0.0,
            "uptime_seconds": 0,
            "uptime_hours": 0.0,
            "uptime_formatted": "unknown",
        })

    return metrics


def _safe_int(value: Any, default: int = 0) -> int:
    try:
        return int(value)
    except (TypeError, ValueError):
        return default


def _safe_float(value: Any, default: float = 0.0) -> float:
    try:
        return float(value)
    except (TypeError, ValueError):
        return default


async def _run_blocking_with_timeout(
    label: str,
    func: Callable[..., Any],
    *args: Any,
    timeout_seconds: float = BLOCKING_CALL_TIMEOUT_SECONDS,
    fallback: Any = None,
) -> Any:
    try:
        return await asyncio.wait_for(asyncio.to_thread(func, *args), timeout=timeout_seconds)
    except asyncio.TimeoutError:
        logger.warning(f"{label} timed out after {timeout_seconds:.2f}s")
    except Exception as exc:
        logger.warning(f"{label} failed: {exc}")
    return fallback


def _should_skip_path(path: Path) -> bool:
    return any(part in IGNORED_SCAN_DIRS for part in path.parts)


def _candidate_urls(env_key: str, defaults: Sequence[str]) -> List[str]:
    values: List[str] = []
    env_value = os.getenv(env_key, "")
    if env_value:
        values.extend(item.strip() for item in env_value.split(",") if item.strip())
    values.extend(defaults)

    deduped: List[str] = []
    for raw in values:
        normalized = raw.rstrip("/")
        if normalized and normalized not in deduped:
            deduped.append(normalized)
    return deduped


def _categorize_material(path: Path) -> Optional[str]:
    suffix = path.suffix.lower()

    if suffix in {".md", ".txt", ".html", ".rst"}:
        return "documentation"
    if suffix in {".json", ".jsonl", ".csv", ".xlsx", ".xls", ".parquet"}:
        return "data"
    if suffix in {".pptx", ".ppt", ".pdf", ".docx"}:
        return "presentations"
    if suffix == ".ipynb":
        return "notebooks"
    return None


def get_project_materials_snapshot(force_refresh: bool = False, limit: int = 12) -> Dict[str, Any]:
    now = time.time()
    cached = MATERIALS_CACHE.get("data")
    if not force_refresh and cached and (now - float(MATERIALS_CACHE.get("fetched_at", 0.0))) < MATERIALS_CACHE_TTL_SECONDS:
        return cached

    extension_counts: Dict[str, int] = {}
    category_counts: Dict[str, int] = {
        "documentation": 0,
        "data": 0,
        "presentations": 0,
        "notebooks": 0,
    }
    recent_materials: List[Dict[str, Any]] = []
    total_materials = 0
    scanned_files = 0
    scan_started = time.perf_counter()
    scan_truncated = False

    if PROJECT_ROOT.exists():
        for current_root, dirnames, filenames in os.walk(PROJECT_ROOT):
            dirnames[:] = [name for name in dirnames if name not in IGNORED_SCAN_DIRS]
            current_path = Path(current_root)

            if _should_skip_path(current_path):
                dirnames[:] = []
                continue

            for filename in filenames:
                scanned_files += 1
                elapsed = time.perf_counter() - scan_started
                if scanned_files > MATERIALS_SCAN_MAX_FILES or elapsed >= MATERIALS_SCAN_MAX_SECONDS:
                    scan_truncated = True
                    break

                file_path = current_path / filename
                category = _categorize_material(file_path)
                if not category:
                    continue

                try:
                    stat = file_path.stat()
                except OSError:
                    continue

                total_materials += 1
                suffix = file_path.suffix.lower() or "<none>"
                category_counts[category] += 1
                extension_counts[suffix] = extension_counts.get(suffix, 0) + 1
                recent_materials.append({
                    "path": file_path.relative_to(PROJECT_ROOT).as_posix(),
                    "category": category,
                    "size_kb": round(stat.st_size / 1024, 1),
                    "modified_at": datetime.fromtimestamp(stat.st_mtime).isoformat(timespec="seconds"),
                    "mtime": stat.st_mtime,
                })

            if scan_truncated:
                break

    recent_materials.sort(key=lambda item: item["mtime"], reverse=True)

    service_dir = PROJECT_ROOT / "services"
    service_samples = []
    if service_dir.exists():
        service_samples = sorted([item.name for item in service_dir.iterdir() if item.is_dir()])

    generated_reports = []
    for report in sorted(REPORTS_DIR.glob("*"), key=lambda item: item.stat().st_mtime, reverse=True)[:limit]:
        report_stat = report.stat()
        generated_reports.append({
            "name": report.name,
            "size_kb": round(report_stat.st_size / 1024, 1),
            "modified_at": datetime.fromtimestamp(report_stat.st_mtime).isoformat(timespec="seconds"),
        })

    snapshot = {
        "root": PROJECT_ROOT.as_posix(),
        "total_materials": total_materials,
        "category_count": sum(1 for count in category_counts.values() if count > 0),
        "categories": category_counts,
        "top_extensions": [
            {"extension": ext, "count": count}
            for ext, count in sorted(extension_counts.items(), key=lambda item: item[1], reverse=True)[:limit]
        ],
        "recent_materials": [
            {key: value for key, value in item.items() if key != "mtime"}
            for item in recent_materials[:limit]
        ],
        "service_directories": len(service_samples),
        "service_samples": service_samples[:limit],
        "generated_reports": generated_reports,
        "scan_truncated": scan_truncated,
        "scanned_files": scanned_files,
        "scan_duration_ms": round((time.perf_counter() - scan_started) * 1000, 1),
    }

    MATERIALS_CACHE["fetched_at"] = now
    MATERIALS_CACHE["data"] = snapshot
    return snapshot


async def _probe_candidates(
    client: httpx.AsyncClient,
    *,
    name: str,
    label: str,
    urls: Sequence[str],
    path: str,
    headers: Optional[Dict[str, str]] = None,
) -> Dict[str, Any]:
    last_error = None
    for base_url in urls:
        started_at = time.perf_counter()
        try:
            response = await client.get(f"{base_url}{path}", headers=headers)
            latency_ms = round((time.perf_counter() - started_at) * 1000, 2)
            payload: Any
            content_type = (response.headers.get("content-type") or "").lower()
            if "json" in content_type:
                payload = response.json()
            else:
                payload = {"text": response.text[:300]}

            summary = payload if isinstance(payload, dict) else {"value": str(payload)[:200]}
            return {
                "name": name,
                "label": label,
                "base_url": base_url,
                "path": path,
                "ok": response.is_success,
                "status_code": response.status_code,
                "latency_ms": latency_ms,
                "payload": payload,
                "summary": {
                    key: summary.get(key)
                    for key in (
                        "status", "service", "requests", "success", "failed", "hit_rate",
                        "redis", "routing_enabled", "queue_depth", "last_request_at"
                    )
                    if isinstance(summary, dict) and key in summary
                },
            }
        except Exception as exc:
            last_error = str(exc)

    return {
        "name": name,
        "label": label,
        "base_url": urls[0] if urls else "",
        "path": path,
        "ok": False,
        "status_code": 0,
        "latency_ms": None,
        "payload": {},
        "summary": {},
        "error": last_error or "unreachable",
    }


async def collect_upstream_snapshot(force_refresh: bool = False) -> Dict[str, Any]:
    now = time.time()
    cached = UPSTREAM_CACHE.get("data")
    if not force_refresh and cached and (now - float(UPSTREAM_CACHE.get("fetched_at", 0.0))) < UPSTREAM_CACHE_TTL_SECONDS:
        return cached

    probe_specs: List[Dict[str, Any]] = [
        {
            "name": "main_api_health",
            "label": "Main API Health",
            "urls": _candidate_urls("CLISONIX_API_URL", ["http://api:8000", "http://localhost:8000"]),
            "path": "/health",
            "headers": None,
        },
        {
            "name": "ocean_document_capabilities",
            "label": "Ocean Document Capabilities",
            "urls": _candidate_urls("OCEAN_CORE_URL", ["http://ocean-core:8030", "http://localhost:8030"]),
            "path": "/api/v1/documents/capabilities",
            "headers": None,
        },
        {
            "name": "ocean_document_metrics",
            "label": "Ocean Document Metrics",
            "urls": _candidate_urls("OCEAN_CORE_URL", ["http://ocean-core:8030", "http://localhost:8030"]),
            "path": "/api/v1/documents/metrics",
            "headers": None,
        },
        {
            "name": "ocean_signal_status",
            "label": "Ocean Signal Status",
            "urls": _candidate_urls("OCEAN_CORE_URL", ["http://ocean-core:8030", "http://localhost:8030"]),
            "path": "/api/v1/signals/status",
            "headers": None,
        },
        {
            "name": "ocean_cache_hit_rate",
            "label": "Ocean Predictive Cache",
            "urls": _candidate_urls("OCEAN_CORE_URL", ["http://ocean-core:8030", "http://localhost:8030"]),
            "path": "/api/v1/v6/cache/hit_rate",
            "headers": None,
        },
        {
            "name": "usage_analytics_status",
            "label": "Usage Analytics Status",
            "urls": _candidate_urls("ANALYTICS_URL", ["http://usage-analytics:8006", "http://localhost:8006"]),
            "path": "/status",
            "headers": None,
        },
    ]

    analytics_bearer = os.getenv("REPORTING_ANALYTICS_BEARER_TOKEN", "").strip()
    if analytics_bearer:
        probe_specs.append({
            "name": "usage_analytics_usage",
            "label": "Usage Analytics Metrics",
            "urls": _candidate_urls("ANALYTICS_URL", ["http://usage-analytics:8006", "http://localhost:8006"]),
            "path": "/api/v1/usage",
            "headers": {"Authorization": f"Bearer {analytics_bearer}"},
        })

    async with httpx.AsyncClient(timeout=UPSTREAM_REQUEST_TIMEOUT_SECONDS, follow_redirects=True) as client:
        results = await asyncio.gather(*[
            _probe_candidates(
                client,
                name=spec["name"],
                label=spec["label"],
                urls=spec["urls"],
                path=spec["path"],
                headers=spec["headers"],
            )
            for spec in probe_specs
        ])

    snapshot = {
        "timestamp": datetime.now().isoformat(),
        "results": results,
        "named": {item["name"]: item for item in results},
        "summary": {
            "total": len(results),
            "healthy": sum(1 for item in results if item.get("ok")),
            "failed": sum(1 for item in results if not item.get("ok")),
        },
    }

    UPSTREAM_CACHE["fetched_at"] = now
    UPSTREAM_CACHE["data"] = snapshot
    return snapshot


def _build_alerts(system: Dict[str, Any], containers: List[Dict[str, Any]], probes: Dict[str, Any]) -> List[Dict[str, Any]]:
    alerts: List[Dict[str, Any]] = []
    timestamp = datetime.now().isoformat()

    cpu_percent = _safe_float(system.get("cpu_percent"))
    memory_percent = _safe_float(system.get("memory_percent"))
    disk_percent = _safe_float(system.get("disk_percent"))

    if cpu_percent >= 85:
        alerts.append({
            "severity": "critical" if cpu_percent >= 95 else "warning",
            "name": "HighCpuUsage",
            "message": f"CPU usage is {cpu_percent:.1f}%",
            "fired_at": timestamp,
            "value": cpu_percent,
        })
    if memory_percent >= 85:
        alerts.append({
            "severity": "critical" if memory_percent >= 92 else "warning",
            "name": "HighMemoryUsage",
            "message": f"Memory usage is {memory_percent:.1f}%",
            "fired_at": timestamp,
            "value": memory_percent,
        })
    if disk_percent >= 90:
        alerts.append({
            "severity": "critical" if disk_percent >= 95 else "warning",
            "name": "HighDiskUsage",
            "message": f"Disk usage is {disk_percent:.1f}%",
            "fired_at": timestamp,
            "value": disk_percent,
        })

    unhealthy_containers = [item.get("name", "unknown") for item in containers if not item.get("healthy")]
    if unhealthy_containers:
        alerts.append({
            "severity": "warning",
            "name": "ContainerAttentionNeeded",
            "message": f"{len(unhealthy_containers)} containers need attention",
            "fired_at": timestamp,
            "value": unhealthy_containers[:5],
        })

    for probe in probes.get("results", []):
        if not probe.get("ok"):
            alerts.append({
                "severity": "warning",
                "name": f"Probe:{probe.get('name', 'unknown')}",
                "message": probe.get("error") or f"{probe.get('label', 'service')} returned {probe.get('status_code', 0)}",
                "fired_at": timestamp,
                "value": probe.get("status_code", 0),
            })

    signal_payload = probes.get("named", {}).get("ocean_signal_status", {}).get("payload", {})
    queue_depth = _safe_int(signal_payload.get("queue_depth"))
    if queue_depth > 0:
        alerts.append({
            "severity": "info",
            "name": "SignalQueueActivity",
            "message": f"Signal queue currently holds {queue_depth} items",
            "fired_at": timestamp,
            "value": queue_depth,
        })

    severity_rank = {"critical": 0, "warning": 1, "info": 2}
    alerts.sort(key=lambda item: severity_rank.get(str(item.get("severity", "info")).lower(), 3))
    return alerts


def _record_metrics_snapshot(system: Dict[str, Any], dashboard: Dict[str, Any]) -> None:
    snapshot = {
        "recorded_at": time.time(),
        "timestamp": dashboard.get("timestamp", datetime.now().isoformat()),
        "system": {
            "cpu_percent": _safe_float(system.get("cpu_percent")),
            "memory_percent": _safe_float(system.get("memory_percent")),
            "disk_percent": _safe_float(system.get("disk_percent")),
            "uptime_seconds": _safe_int(system.get("uptime_seconds")),
        },
        "api_requests_24h": _safe_int(dashboard.get("api_requests_24h")),
        "api_errors_24h": _safe_int(dashboard.get("api_errors_24h")),
        "documents_generated": _safe_int(dashboard.get("documents_generated")),
        "cache_hit_rate": _safe_float(dashboard.get("cache_hit_rate")),
    }
    METRIC_HISTORY.append(snapshot)
    if len(METRIC_HISTORY) > MAX_HISTORY_POINTS:
        del METRIC_HISTORY[:-MAX_HISTORY_POINTS]


async def build_dashboard_payload() -> Dict[str, Any]:
    system, containers, stats, probes, materials = await asyncio.gather(
        _run_blocking_with_timeout(
            "system_metrics",
            get_system_metrics_real,
            timeout_seconds=0.9,
            fallback={
                "cpu_percent": 0.0,
                "memory_percent": 0.0,
                "disk_percent": 0.0,
                "uptime_seconds": 0,
                "uptime_formatted": "unknown",
            },
        ),
        _run_blocking_with_timeout("docker_containers", get_docker_containers_real, timeout_seconds=1.25, fallback=[]),
        _run_blocking_with_timeout("docker_stats", get_docker_stats_real, timeout_seconds=1.25, fallback=[]),
        collect_upstream_snapshot(),
        _run_blocking_with_timeout(
            "project_materials",
            get_project_materials_snapshot,
            timeout_seconds=1.0,
            fallback={"total_materials": 0, "category_count": 0, "generated_reports": []},
        ),
    )
    apis = get_api_endpoints_real()

    doc_metrics = probes.get("named", {}).get("ocean_document_metrics", {}).get("payload", {}) or {}
    cache_metrics = probes.get("named", {}).get("ocean_cache_hit_rate", {}).get("payload", {}) or {}
    usage_metrics = probes.get("named", {}).get("usage_analytics_usage", {}).get("payload", {}) or {}

    api_requests_24h = _safe_int(usage_metrics.get("requests_today"), _safe_int(doc_metrics.get("requests")))
    api_errors_24h = _safe_int(doc_metrics.get("failed"))
    documents_generated = _safe_int(doc_metrics.get("success"), _safe_int(doc_metrics.get("requests")))
    cache_hit_rate = round(_safe_float(cache_metrics.get("hit_rate")) * 100, 2)

    successful_latencies = [
        _safe_float(item.get("latency_ms"))
        for item in probes.get("results", [])
        if item.get("ok") and item.get("latency_ms") is not None
    ]
    avg_response_ms = round(sum(successful_latencies) / len(successful_latencies), 2) if successful_latencies else 0.0

    alerts = _build_alerts(system, containers, probes)
    healthy_containers = sum(1 for item in containers if item.get("healthy"))

    payload = {
        "timestamp": datetime.now().isoformat(),
        "data_type": "REAL",
        "status": "operational" if not alerts else "degraded",
        "system": system,
        "docker": {
            "containers": containers,
            "stats": stats,
            "total": len(containers),
            "healthy": healthy_containers,
        },
        "api_catalog": {
            "total": len(apis),
            "endpoints": apis,
        },
        "api_requests_24h": api_requests_24h,
        "api_errors_24h": api_errors_24h,
        "documents_generated": documents_generated,
        "cache_hit_rate": cache_hit_rate,
        "avg_response_ms": avg_response_ms,
        "service_fleet_total": probes.get("summary", {}).get("total", 0),
        "running_containers": healthy_containers,
        "total_containers": len(containers),
        "alerts": alerts,
        "upstream": probes,
        "project_materials": materials,
        "insights": {
            "document_service_available": bool(probes.get("named", {}).get("ocean_document_capabilities", {}).get("ok")),
            "cache_metrics_available": bool(probes.get("named", {}).get("ocean_cache_hit_rate", {}).get("ok")),
            "usage_analytics_available": bool(probes.get("named", {}).get("usage_analytics_status", {}).get("ok")),
            "generated_reports": len(materials.get("generated_reports", [])),
        },
    }

    _record_metrics_snapshot(system, payload)
    return payload


def get_api_endpoints_real() -> List[Dict[str, Any]]:
    """Lista e API endpoints REALE nga aplikacioni aktual dhe integrimet kryesore."""
    endpoints: List[Dict[str, Any]] = []
    seen = set()

    for route in sorted(app.routes, key=lambda item: getattr(item, "path", "")):
        path = getattr(route, "path", None)
        methods = sorted(
            method for method in (getattr(route, "methods", set()) or set())
            if method not in {"HEAD", "OPTIONS"}
        )
        if not path or not methods:
            continue

        key = (path, tuple(methods))
        if key in seen:
            continue
        seen.add(key)

        folder = path.strip("/").split("/")[0] if path.strip("/") else "root"
        endpoints.append({
            "id": f"RPT-{len(endpoints) + 1:03d}",
            "method": ", ".join(methods),
            "endpoint": path,
            "folder": folder or "root",
            "status": "LOCAL",
        })

    upstream_catalog = [
        ("GET", "/api/v1/documents/capabilities", "Ocean", "UPSTREAM"),
        ("GET", "/api/v1/documents/metrics", "Ocean", "UPSTREAM"),
        ("GET", "/api/v1/signals/status", "Signals", "UPSTREAM"),
        ("GET", "/api/v1/v6/cache/hit_rate", "Ocean", "UPSTREAM"),
        ("GET", "/status", "Analytics", "UPSTREAM"),
    ]
    for method, path, folder, status in upstream_catalog:
        if (path, (method,)) in seen:
            continue
        endpoints.append({
            "id": f"RPT-{len(endpoints) + 1:03d}",
            "method": method,
            "endpoint": path,
            "folder": folder,
            "status": status,
        })

    return endpoints


@app.get("/health")
@app.get("/api/reporting/health")
async def health():
    materials = await _run_blocking_with_timeout(
        "health_materials",
        get_project_materials_snapshot,
        timeout_seconds=1.0,
        fallback={"total_materials": 0},
    )
    return {
        "status": "healthy",
        "service": "reporting-real-client-intelligence",
        "excel_available": EXCEL_AVAILABLE,
        "pptx_available": PPTX_AVAILABLE,
        "psutil_available": PSUTIL_AVAILABLE,
        "docker_available": bool(get_docker_client()),
        "project_root": PROJECT_ROOT.as_posix(),
        "material_count": materials.get("total_materials", 0),
        "features": [
            "real_system_metrics",
            "docker_reporting",
            "document_capabilities",
            "project_material_inventory",
            "metrics_history",
            "excel_export",
            "powerpoint_export",
        ],
        "timestamp": datetime.now().isoformat(),
    }


@app.get("/status")
@app.get("/api/status")
async def status():
    """Status endpoint production-grade me real summaries."""
    containers, system, upstream = await asyncio.gather(
        _run_blocking_with_timeout("status_containers", get_docker_containers_real, timeout_seconds=1.25, fallback=[]),
        _run_blocking_with_timeout(
            "status_system",
            get_system_metrics_real,
            timeout_seconds=0.9,
            fallback={"cpu_percent": 0.0, "memory_percent": 0.0, "disk_percent": 0.0, "uptime_seconds": 0, "uptime_formatted": "unknown"},
        ),
        collect_upstream_snapshot(),
    )

    return {
        "status": "active" if upstream.get("summary", {}).get("failed", 0) == 0 else "degraded",
        "service": "reporting-microservice",
        "uptime": system.get("uptime_formatted", "unknown"),
        "timestamp": datetime.now().isoformat(),
        "containers": {
            "total": len(containers),
            "healthy": sum(1 for c in containers if c.get("healthy")),
        },
        "system": {
            "cpu_percent": system.get("cpu_percent", 0),
            "memory_percent": system.get("memory_percent", 0),
            "disk_percent": system.get("disk_percent", 0),
            "uptime_seconds": system.get("uptime_seconds", 0),
        },
        "upstream": upstream.get("summary", {}),
    }


@app.get("/api/reporting/docker-containers")
async def docker_containers():
    """Merr listën e Docker containers REALE"""
    containers = await _run_blocking_with_timeout(
        "endpoint_docker_containers",
        get_docker_containers_real,
        timeout_seconds=1.25,
        fallback=[],
    )
    return {
        "timestamp": datetime.now().isoformat(),
        "data_type": "REAL",
        "total": len(containers),
        "healthy": sum(1 for c in containers if c.get("healthy")),
        "running": sum(1 for c in containers if c.get("healthy")),
        "containers": containers,
    }


@app.get("/api/reporting/docker-stats")
async def docker_stats():
    """Merr CPU/Memory stats REALE për Docker containers"""
    stats = await _run_blocking_with_timeout(
        "endpoint_docker_stats",
        get_docker_stats_real,
        timeout_seconds=1.25,
        fallback=[],
    )
    return {
        "timestamp": datetime.now().isoformat(),
        "data_type": "REAL",
        "total": len(stats),
        "stats": stats,
    }


@app.get("/api/reporting/system-metrics")
async def system_metrics():
    """Merr metrika REALE nga sistemi (CPU, RAM, Disk)"""
    metrics = await _run_blocking_with_timeout(
        "endpoint_system_metrics",
        get_system_metrics_real,
        timeout_seconds=0.9,
        fallback={"cpu_percent": 0.0, "memory_percent": 0.0, "disk_percent": 0.0, "uptime_seconds": 0},
    )
    return {
        "timestamp": datetime.now().isoformat(),
        "data_type": "REAL",
        "metrics": metrics,
    }


@app.get("/api/reporting/capabilities")
async def reporting_capabilities():
    materials = get_project_materials_snapshot()
    upstream = await collect_upstream_snapshot()
    document_capabilities = upstream.get("named", {}).get("ocean_document_capabilities", {}).get("payload", {})

    export_formats = ["xlsx"] + (["pptx"] if PPTX_AVAILABLE else [])

    return {
        "timestamp": datetime.now().isoformat(),
        "service": "reporting-real-client-intelligence",
        "status": "operational" if PPTX_AVAILABLE else "degraded",
        "exports": export_formats,
        "real_data_sources": [
            "system metrics via psutil",
            "docker runtime inspection",
            "Curiosity Ocean document capabilities",
            "signal routing status",
            "predictive cache metrics",
            "project materials inventory",
        ],
        "documents": document_capabilities,
        "project_materials": {
            "total_materials": materials.get("total_materials", 0),
            "category_count": materials.get("category_count", 0),
            "service_directories": materials.get("service_directories", 0),
        },
        "api_catalog_total": len(get_api_endpoints_real()),
    }


@app.get("/api/reporting/project-materials")
async def project_materials(force_refresh: bool = Query(False)):
    """Inventory real i dokumenteve, datasets, notebooks dhe materialeve të projektit."""
    return {
        "timestamp": datetime.now().isoformat(),
        "data_type": "REAL",
        "materials": get_project_materials_snapshot(force_refresh=force_refresh),
    }


@app.get("/api/reporting/metrics-history")
async def metrics_history(hours: int = Query(24, ge=1, le=720)):
    """Historik real i snapshot-eve të metrikave gjatë jetës së procesit."""
    if not METRIC_HISTORY:
        await build_dashboard_payload()

    cutoff = time.time() - (hours * 3600)
    history = [item for item in METRIC_HISTORY if float(item.get("recorded_at", 0.0)) >= cutoff]
    return {
        "timestamp": datetime.now().isoformat(),
        "data_type": "REAL",
        "period_hours": hours,
        "data_points": len(history),
        "history": history,
    }


@app.get("/api/reporting/alerts")
async def reporting_alerts():
    payload = await build_dashboard_payload()
    return {
        "timestamp": payload.get("timestamp"),
        "status": payload.get("status"),
        "alerts": payload.get("alerts", []),
    }


@app.get("/api/reporting/list-reports")
async def list_reports():
    items = []
    for report in sorted(REPORTS_DIR.glob("*"), key=lambda item: item.stat().st_mtime, reverse=True):
        if not report.is_file():
            continue
        stat = report.stat()
        items.append({
            "name": report.name,
            "extension": report.suffix.lower(),
            "size_bytes": stat.st_size,
            "size_kb": round(stat.st_size / 1024, 1),
            "modified_at": datetime.fromtimestamp(stat.st_mtime).isoformat(timespec="seconds"),
            "download_url": f"/api/reporting/download/{report.name}",
        })

    return {
        "timestamp": datetime.now().isoformat(),
        "total": len(items),
        "reports": items,
    }


@app.get("/api/reporting/download/{filename}")
async def download_report(filename: str):
    filepath = (REPORTS_DIR / filename).resolve()
    if REPORTS_DIR.resolve() not in filepath.parents or not filepath.exists() or not filepath.is_file():
        raise HTTPException(status_code=404, detail="Report not found")
    return FileResponse(path=filepath, filename=filepath.name)


@app.delete("/api/reporting/clear-reports")
async def clear_reports(older_than_hours: int = Query(168, ge=0, le=8760)):
    cutoff = time.time() - (older_than_hours * 3600)
    deleted = []

    for report in REPORTS_DIR.glob("*"):
        if not report.is_file():
            continue
        if older_than_hours == 0 or report.stat().st_mtime < cutoff:
            deleted.append(report.name)
            report.unlink(missing_ok=True)

    MATERIALS_CACHE["fetched_at"] = 0.0
    MATERIALS_CACHE["data"] = None

    return {
        "timestamp": datetime.now().isoformat(),
        "deleted_count": len(deleted),
        "deleted": deleted,
        "older_than_hours": older_than_hours,
    }


@app.get("/api/reporting/dashboard")
async def dashboard():
    """Dashboard me të dhëna 100% REALE dhe klient-ready insights."""
    return await build_dashboard_payload()


@app.get("/api/reporting/export-excel")
async def export_excel():
    """
    Gjeneron Excel REAL me tabela të plota për operacione, dokumente dhe analytics klientësh.
    """
    dependencies = [Workbook, Table, TableStyleInfo, Font, PatternFill, Alignment, Border, Side, get_column_letter]
    if not EXCEL_AVAILABLE or any(dep is None for dep in dependencies):
        raise HTTPException(status_code=500, detail="openpyxl components not available")

    try:
        containers = get_docker_containers_real()
        stats = get_docker_stats_real()
        system = get_system_metrics_real()
        apis = get_api_endpoints_real()
        probes = await collect_upstream_snapshot()
        materials = get_project_materials_snapshot()

        assert Workbook is not None and Table is not None and TableStyleInfo is not None
        assert Font is not None and PatternFill is not None and Alignment is not None
        assert Border is not None and Side is not None and get_column_letter is not None

        workbook_cls = cast(Any, Workbook)
        table_cls = cast(Any, Table)
        table_style_cls = cast(Any, TableStyleInfo)
        font_cls = cast(Any, Font)
        fill_cls = cast(Any, PatternFill)
        alignment_cls = cast(Any, Alignment)
        border_cls = cast(Any, Border)
        side_cls = cast(Any, Side)
        get_col_letter = cast(Any, get_column_letter)

        wb = workbook_cls()

        header_font = font_cls(bold=True, color="FFFFFF", size=11)
        header_fill = fill_cls(start_color="1F4E79", end_color="1F4E79", fill_type="solid")
        ready_fill = fill_cls(start_color="C6EFCE", end_color="C6EFCE", fill_type="solid")
        border = border_cls(
            left=side_cls(style='thin'), right=side_cls(style='thin'),
            top=side_cls(style='thin'), bottom=side_cls(style='thin')
        )
        table_style = table_style_cls(
            name="TableStyleMedium2", showFirstColumn=False,
            showLastColumn=False, showRowStripes=True, showColumnStripes=False
        )

        ws1 = wb.active
        ws1.title = "Docker Containers"
        headers = ["Container ID", "Name", "Image", "Status", "Ports", "Health", "Uptime"]
        for col, header in enumerate(headers, 1):
            cell = ws1.cell(row=1, column=col, value=header)
            cell.font = header_font
            cell.fill = header_fill
            cell.border = border
            cell.alignment = alignment_cls(horizontal="center")

        for row_num, container in enumerate(containers, 2):
            ws1.cell(row=row_num, column=1, value=container.get("container_id", "")).border = border
            ws1.cell(row=row_num, column=2, value=container.get("name", "")).border = border
            ws1.cell(row=row_num, column=3, value=container.get("image", "")).border = border
            ws1.cell(row=row_num, column=4, value=container.get("status", "")).border = border
            ws1.cell(row=row_num, column=5, value=container.get("ports", "")).border = border
            health_cell = ws1.cell(row=row_num, column=6, value="Healthy" if container.get("healthy") else "Attention")
            health_cell.border = border
            if container.get("healthy"):
                health_cell.fill = ready_fill
            ws1.cell(row=row_num, column=7, value=container.get("uptime", "")).border = border

        if containers:
            table = table_cls(displayName="DockerContainers", ref=f"A1:G{len(containers) + 1}")
            table.tableStyleInfo = table_style
            ws1.add_table(table)

        for col in range(1, 8):
            ws1.column_dimensions[get_col_letter(col)].width = 25

        ws2 = wb.create_sheet("Container Stats")
        headers2 = ["Container", "CPU %", "Memory Usage", "Memory %", "Network I/O", "Block I/O"]
        for col, header in enumerate(headers2, 1):
            cell = ws2.cell(row=1, column=col, value=header)
            cell.font = header_font
            cell.fill = header_fill
            cell.border = border

        for row_num, stat in enumerate(stats, 2):
            ws2.cell(row=row_num, column=1, value=stat.get("container", "")).border = border
            ws2.cell(row=row_num, column=2, value=stat.get("cpu", "")).border = border
            ws2.cell(row=row_num, column=3, value=stat.get("mem_usage", "")).border = border
            ws2.cell(row=row_num, column=4, value=stat.get("mem_percent", "")).border = border
            ws2.cell(row=row_num, column=5, value=stat.get("net_io", "")).border = border
            ws2.cell(row=row_num, column=6, value=stat.get("block_io", "")).border = border

        if stats:
            table2 = table_cls(displayName="ContainerStats", ref=f"A1:F{len(stats) + 1}")
            table2.tableStyleInfo = table_style
            ws2.add_table(table2)

        for col in range(1, 7):
            ws2.column_dimensions[get_col_letter(col)].width = 22

        ws3 = wb.create_sheet("System Metrics")
        headers3 = ["Metric", "Value", "Unit", "Status"]
        for col, header in enumerate(headers3, 1):
            cell = ws3.cell(row=1, column=col, value=header)
            cell.font = header_font
            cell.fill = header_fill
            cell.border = border

        system_data = [
            ("CPU Usage", system.get("cpu_percent", 0), "%", "Normal" if system.get("cpu_percent", 0) < 80 else "High"),
            ("CPU Cores", system.get("cpu_count", 0), "cores", "-"),
            ("Memory Used", system.get("memory_used_gb", 0), "GB", "Normal" if system.get("memory_percent", 0) < 85 else "High"),
            ("Memory Total", system.get("memory_total_gb", 0), "GB", "-"),
            ("Memory %", system.get("memory_percent", 0), "%", "Normal" if system.get("memory_percent", 0) < 85 else "High"),
            ("Disk Used", system.get("disk_used_gb", 0), "GB", "Normal" if system.get("disk_percent", 0) < 90 else "Warning"),
            ("Disk Total", system.get("disk_total_gb", 0), "GB", "-"),
            ("Disk %", system.get("disk_percent", 0), "%", "Normal" if system.get("disk_percent", 0) < 90 else "Warning"),
            ("Network Sent", system.get("net_sent_gb", 0), "GB", "-"),
            ("Network Received", system.get("net_recv_gb", 0), "GB", "-"),
            ("System Uptime", system.get("uptime_formatted", "unknown"), "formatted", "-"),
        ]

        for row_num, (metric, value, unit, status) in enumerate(system_data, 2):
            ws3.cell(row=row_num, column=1, value=metric).border = border
            ws3.cell(row=row_num, column=2, value=value).border = border
            ws3.cell(row=row_num, column=3, value=unit).border = border
            status_cell = ws3.cell(row=row_num, column=4, value=status)
            status_cell.border = border
            if status == "Normal":
                status_cell.fill = ready_fill

        table3 = table_cls(displayName="SystemMetrics", ref=f"A1:D{len(system_data) + 1}")
        table3.tableStyleInfo = table_style
        ws3.add_table(table3)
        for col in range(1, 5):
            ws3.column_dimensions[get_col_letter(col)].width = 20

        ws4 = wb.create_sheet("API Endpoints")
        headers4 = ["ID", "Method", "Endpoint", "Folder", "Status"]
        for col, header in enumerate(headers4, 1):
            cell = ws4.cell(row=1, column=col, value=header)
            cell.font = header_font
            cell.fill = header_fill
            cell.border = border

        for row_num, api in enumerate(apis, 2):
            ws4.cell(row=row_num, column=1, value=api.get("id", "")).border = border
            ws4.cell(row=row_num, column=2, value=api.get("method", "")).border = border
            ws4.cell(row=row_num, column=3, value=api.get("endpoint", "")).border = border
            ws4.cell(row=row_num, column=4, value=api.get("folder", "")).border = border
            status_cell = ws4.cell(row=row_num, column=5, value=api.get("status", ""))
            status_cell.border = border
            if api.get("status") in {"LOCAL", "UPSTREAM"}:
                status_cell.fill = ready_fill

        if apis:
            table4 = table_cls(displayName="APIEndpoints", ref=f"A1:E{len(apis) + 1}")
            table4.tableStyleInfo = table_style
            ws4.add_table(table4)
        for col in range(1, 6):
            ws4.column_dimensions[get_col_letter(col)].width = 28

        ws5 = wb.create_sheet("Service Probes")
        headers5 = ["Service", "URL", "Status Code", "OK", "Latency (ms)", "Summary"]
        for col, header in enumerate(headers5, 1):
            cell = ws5.cell(row=1, column=col, value=header)
            cell.font = header_font
            cell.fill = header_fill
            cell.border = border

        probe_results = probes.get("results", [])
        for row_num, probe in enumerate(probe_results, 2):
            ws5.cell(row=row_num, column=1, value=probe.get("label", probe.get("name", ""))).border = border
            ws5.cell(row=row_num, column=2, value=f"{probe.get('base_url', '')}{probe.get('path', '')}").border = border
            ws5.cell(row=row_num, column=3, value=probe.get("status_code", 0)).border = border
            ok_cell = ws5.cell(row=row_num, column=4, value="Yes" if probe.get("ok") else "No")
            ok_cell.border = border
            if probe.get("ok"):
                ok_cell.fill = ready_fill
            ws5.cell(row=row_num, column=5, value=probe.get("latency_ms") or "-").border = border
            ws5.cell(row=row_num, column=6, value=json.dumps(probe.get("summary", {}), ensure_ascii=False)[:300]).border = border

        if probe_results:
            table5 = table_cls(displayName="ServiceProbes", ref=f"A1:F{len(probe_results) + 1}")
            table5.tableStyleInfo = table_style
            ws5.add_table(table5)
        for col in range(1, 7):
            ws5.column_dimensions[get_col_letter(col)].width = 26

        ws6 = wb.create_sheet("Project Materials")
        headers6 = ["Section", "Name", "Value", "Notes"]
        for col, header in enumerate(headers6, 1):
            cell = ws6.cell(row=1, column=col, value=header)
            cell.font = header_font
            cell.fill = header_fill
            cell.border = border

        materials_rows = [
            ("Inventory", "Total materials", materials.get("total_materials", 0), PROJECT_ROOT.as_posix()),
            ("Inventory", "Category count", materials.get("category_count", 0), "documentation/data/presentations/notebooks"),
            ("Inventory", "Service directories", materials.get("service_directories", 0), ", ".join(materials.get("service_samples", [])[:5])),
        ]
        for name, count in materials.get("categories", {}).items():
            materials_rows.append(("Category", name, count, "real files discovered"))
        for ext_item in materials.get("top_extensions", [])[:8]:
            materials_rows.append(("Extension", ext_item.get("extension", ""), ext_item.get("count", 0), "top project material types"))

        for row_num, (section, name, value, notes) in enumerate(materials_rows, 2):
            ws6.cell(row=row_num, column=1, value=section).border = border
            ws6.cell(row=row_num, column=2, value=name).border = border
            ws6.cell(row=row_num, column=3, value=value).border = border
            ws6.cell(row=row_num, column=4, value=notes).border = border

        if materials_rows:
            table6 = table_cls(displayName="ProjectMaterials", ref=f"A1:D{len(materials_rows) + 1}")
            table6.tableStyleInfo = table_style
            ws6.add_table(table6)
        for col in range(1, 5):
            ws6.column_dimensions[get_col_letter(col)].width = 28

        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
        filename = f"clisonix_real_report_{timestamp}.xlsx"
        filepath = REPORTS_DIR / filename
        wb.save(str(filepath))

        with open(filepath, 'rb') as f:
            content = f.read()

        logger.info(f"Excel REAL generated: {filename}, size: {len(content)} bytes")
        return Response(
            content=content,
            media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            headers={"Content-Disposition": f"attachment; filename={filename}"}
        )

    except Exception as e:
        logger.error(f"Excel error: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/api/reporting/export-pptx")
async def export_pptx():
    """Gjeneron PowerPoint me të dhëna REALE dhe summary klient-ready."""
    try:
        from pptx import Presentation  # type: ignore[import-not-found]
        from pptx.dml.color import RGBColor  # type: ignore[import-not-found]
        from pptx.enum.text import PP_ALIGN  # type: ignore[import-not-found]
        from pptx.util import Inches, Pt  # type: ignore[import-not-found]
    except ImportError:
        raise HTTPException(status_code=503, detail="python-pptx not installed; PPTX export is temporarily unavailable")

    try:
        dashboard_payload = await build_dashboard_payload()
        system = dashboard_payload.get("system", {})
        containers = dashboard_payload.get("docker", {}).get("containers", [])
        alerts = dashboard_payload.get("alerts", [])
        materials = dashboard_payload.get("project_materials", {})

        prs = Presentation()

        slide1 = prs.slides.add_slide(prs.slide_layouts[6])
        title = slide1.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(9), Inches(1))
        tf = title.text_frame
        tf.text = "Clisonix Client Intelligence Report"
        tf.paragraphs[0].font.size = Pt(34)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = RGBColor(31, 78, 121)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        subtitle = slide1.shapes.add_textbox(Inches(0.5), Inches(3.3), Inches(9), Inches(0.6))
        tf2 = subtitle.text_frame
        tf2.text = f"Real operational + document metrics • {datetime.now().strftime('%Y-%m-%d %H:%M')}"
        tf2.paragraphs[0].font.size = Pt(16)
        tf2.paragraphs[0].alignment = PP_ALIGN.CENTER

        slide2 = prs.slides.add_slide(prs.slide_layouts[6])
        title2 = slide2.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9), Inches(0.6))
        title2.text_frame.text = "Production Summary"
        title2.text_frame.paragraphs[0].font.size = Pt(28)
        title2.text_frame.paragraphs[0].font.bold = True

        metrics_text = f"""
CPU: {system.get('cpu_percent', 0)}%
Memory: {system.get('memory_percent', 0)}% ({system.get('memory_used_gb', 0)} / {system.get('memory_total_gb', 0)} GB)
Disk: {system.get('disk_percent', 0)}% ({system.get('disk_used_gb', 0)} / {system.get('disk_total_gb', 0)} GB)
Uptime: {system.get('uptime_formatted', 'unknown')}
Documents generated: {dashboard_payload.get('documents_generated', 0)}
API requests (24h view): {dashboard_payload.get('api_requests_24h', 0)}
Predictive cache hit rate: {dashboard_payload.get('cache_hit_rate', 0)}%
Tracked project materials: {materials.get('total_materials', 0)}
"""

        body = slide2.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(4.6))
        tf3 = body.text_frame
        tf3.text = metrics_text.strip()
        tf3.paragraphs[0].font.size = Pt(18)

        slide3 = prs.slides.add_slide(prs.slide_layouts[6])
        title3 = slide3.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9), Inches(0.6))
        title3.text_frame.text = "Service & Container Status"
        title3.text_frame.paragraphs[0].font.size = Pt(28)
        title3.text_frame.paragraphs[0].font.bold = True

        y_pos = 1.0
        for container in containers[:8]:
            box = slide3.shapes.add_textbox(Inches(0.5), Inches(y_pos), Inches(9), Inches(0.38))
            box_tf = box.text_frame
            status_icon = "✓" if container.get("healthy") else "•"
            box_tf.text = f"{status_icon} {container.get('name')}: {container.get('status')}"
            box_tf.paragraphs[0].font.size = Pt(13)
            if container.get("healthy"):
                box_tf.paragraphs[0].font.color.rgb = RGBColor(0, 128, 0)
            y_pos += 0.42

        slide4 = prs.slides.add_slide(prs.slide_layouts[6])
        title4 = slide4.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9), Inches(0.6))
        title4.text_frame.text = "Operational Alerts"
        title4.text_frame.paragraphs[0].font.size = Pt(28)
        title4.text_frame.paragraphs[0].font.bold = True

        y_pos = 1.0
        for alert in alerts[:8]:
            box = slide4.shapes.add_textbox(Inches(0.5), Inches(y_pos), Inches(9), Inches(0.45))
            box_tf = box.text_frame
            severity = str(alert.get("severity", "info")).upper()
            box_tf.text = f"[{severity}] {alert.get('name')}: {alert.get('message')}"
            box_tf.paragraphs[0].font.size = Pt(12)
            if severity == "CRITICAL":
                box_tf.paragraphs[0].font.color.rgb = RGBColor(192, 0, 0)
            elif severity == "WARNING":
                box_tf.paragraphs[0].font.color.rgb = RGBColor(192, 96, 0)
            else:
                box_tf.paragraphs[0].font.color.rgb = RGBColor(31, 78, 121)
            y_pos += 0.48

        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
        filename = f"clisonix_real_{timestamp}.pptx"
        filepath = REPORTS_DIR / filename
        prs.save(str(filepath))

        with open(filepath, 'rb') as f:
            content = f.read()

        logger.info(f"PPTX REAL generated: {filename}")
        return Response(
            content=content,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={"Content-Disposition": f"attachment; filename={filename}"}
        )

    except Exception as e:
        logger.error(f"PPTX error: {e}")
        raise HTTPException(status_code=500, detail=str(e))

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8001)
